#!/usr/bin/env python3
"""
svg_to_slides.py — Convert SVG files to PPTX with editable DrawingML vector paths.
Each SVG becomes a one-slide PPTX saved alongside the original file, ready to
import into Google Slides (File → Import slides).

Install deps: pip3 install --break-system-packages python-pptx lxml
Usage:        python3 ~/Documents/svg_to_slides.py file.svg [more.svg ...]
"""
import sys, re, math
from pathlib import Path
from xml.etree import ElementTree as ET

_lib = Path.home() / '.local' / 'share' / 'svg-to-slides' / 'lib'
if _lib.exists() and str(_lib) not in sys.path:
    sys.path.insert(0, str(_lib))

try:
    from pptx import Presentation
    from pptx.util import Emu
    from lxml import etree
except ImportError:
    sys.exit("Run Install.command, or: pip3 install python-pptx lxml")

# ── Constants ────────────────────────────────────────────────────────────────

SLIDE_W = 9144000
SLIDE_H = 6858000
COORD   = 100000
KAPPA   = 0.5522847498

_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ── Color utilities ──────────────────────────────────────────────────────────

CSS_COLORS = {
    'black':'#000000','white':'#ffffff','red':'#ff0000','green':'#008000',
    'blue':'#0000ff','yellow':'#ffff00','orange':'#ffa500','purple':'#800080',
    'pink':'#ffc0cb','gray':'#808080','grey':'#808080','brown':'#a52a2a',
    'cyan':'#00ffff','magenta':'#ff00ff','lime':'#00ff00','navy':'#000080',
    'teal':'#008080','silver':'#c0c0c0','maroon':'#800000','olive':'#808000',
    'aqua':'#00ffff','fuchsia':'#ff00ff','coral':'#ff7f50','gold':'#ffd700',
    'indigo':'#4b0082','violet':'#ee82ee','turquoise':'#40e0d0',
    'darkblue':'#00008b','darkgreen':'#006400','darkred':'#8b0000',
    'darkgray':'#a9a9a9','darkgrey':'#a9a9a9','lightgray':'#d3d3d3',
    'lightgrey':'#d3d3d3','lightblue':'#add8e6','beige':'#f5f5dc',
    'ivory':'#fffff0','khaki':'#f0e68c','lavender':'#e6e6fa',
    'lightgreen':'#90ee90','lightyellow':'#ffffe0','limegreen':'#32cd32',
    'orangered':'#ff4500','royalblue':'#4169e1','sienna':'#a0522d',
    'skyblue':'#87ceeb','slategray':'#708090','steelblue':'#4682b4',
    'tan':'#d2b48c','tomato':'#ff6347','wheat':'#f5deb3',
    'whitesmoke':'#f5f5f5','yellowgreen':'#9acd32','crimson':'#dc143c',
    'darkorange':'#ff8c00','deeppink':'#ff1493','dodgerblue':'#1e90ff',
    'firebrick':'#b22222','forestgreen':'#228b22','goldenrod':'#daa520',
    'hotpink':'#ff69b4','mediumblue':'#0000cd','mediumpurple':'#9370db',
    'midnightblue':'#191970','peru':'#cd853f','plum':'#dda0dd',
    'rosybrown':'#bc8f8f','saddlebrown':'#8b4513','salmon':'#fa8072',
    'seagreen':'#2e8b57','snow':'#fffafa','slateblue':'#6a5acd',
    'springgreen':'#00ff7f','darkviolet':'#9400d3','deepskyblue':'#00bfff',
}

def parse_color(s):
    """Parse any CSS color to 6-char uppercase hex. Returns None for none/transparent."""
    if not s: return None
    s = s.strip()
    sl = s.lower()
    if sl in ('none', 'transparent', ''): return None
    if sl == 'currentcolor': return None
    if sl in CSS_COLORS: return CSS_COLORS[sl].lstrip('#').upper()
    if s.startswith('#'):
        h = s[1:]
        if len(h) == 3: h = h[0]*2 + h[1]*2 + h[2]*2
        return h.upper() if len(h) == 6 else None
    m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', sl)
    if m: return '%02X%02X%02X' % (int(m.group(1)), int(m.group(2)), int(m.group(3)))
    m = re.match(r'rgba?\(\s*(\d+(?:\.\d+)?)%\s*,\s*(\d+(?:\.\d+)?)%\s*,\s*(\d+(?:\.\d+)?)%', sl)
    if m:
        def p(v): return round(float(v) / 100 * 255)
        return '%02X%02X%02X' % (p(m.group(1)), p(m.group(2)), p(m.group(3)))
    return '000000'

# ── Gradient resolution ──────────────────────────────────────────────────────

_NS_RE = re.compile(r'\{[^}]+\}')
def _tag(el): return _NS_RE.sub('', el.tag)

def extract_gradients(root):
    """Return {id: hex_color} mapping gradient IDs to their dominant stop color."""
    grads = {}
    for el in root.iter():
        if _tag(el) not in ('linearGradient', 'radialGradient'): continue
        gid = el.get('id')
        if not gid: continue
        stops = []
        for child in el.iter():
            if _tag(child) != 'stop': continue
            style = child.get('style', '')
            m = re.search(r'stop-color\s*:\s*([^;]+)', style)
            raw = m.group(1).strip() if m else child.get('stop-color', '')
            om = re.search(r'stop-opacity\s*:\s*([^;]+)', style)
            opacity = float(om.group(1).strip() if om else child.get('stop-opacity', '1'))
            c = parse_color(raw)
            if c and opacity > 0.1: stops.append(c)
        if stops: grads[gid] = stops[len(stops) // 2]
    return grads

# ── Style parsing ────────────────────────────────────────────────────────────

def parse_style(el):
    props = {}
    for part in el.get('style', '').split(';'):
        if ':' in part:
            k, _, v = part.partition(':')
            props[k.strip().lower()] = v.strip()
    return props

def get_prop(style_dict, el, *keys):
    for k in keys:
        if k in style_dict: return style_dict[k]
        v = el.get(k)
        if v is not None: return v
    return None

def resolve_fill_stroke(el, gradients, inh_fill):
    """Return (fill_hex_or_None, stroke_hex_or_None, stroke_width_float)."""
    style = parse_style(el)
    fill_raw   = get_prop(style, el, 'fill')
    stroke_raw = get_prop(style, el, 'stroke')
    sw_raw     = get_prop(style, el, 'stroke-width')

    if fill_raw and fill_raw.lower().startswith('url(#'):
        m = re.match(r'url\(#([^)]+)\)', fill_raw)
        fill_hex = gradients.get(m.group(1)) if m else None
        if fill_hex is None: fill_hex = inh_fill
    elif fill_raw:
        fill_hex = parse_color(fill_raw)
    else:
        fill_hex = inh_fill

    stroke_hex = parse_color(stroke_raw) if stroke_raw else None
    try:
        sw = float(re.sub(r'[^\d.]', '', sw_raw)) if sw_raw else 0.0
    except ValueError:
        sw = 0.0

    return fill_hex, stroke_hex, sw

# ── Basic shapes → path d ────────────────────────────────────────────────────

def _f(*vals):
    return ' '.join(str(round(v, 4)) for v in vals)

def rect_to_d(el):
    x  = float(el.get('x', 0)); y  = float(el.get('y', 0))
    w  = float(el.get('width', 0)); h  = float(el.get('height', 0))
    rx = float(el.get('rx') or el.get('ry') or 0)
    ry = float(el.get('ry') or el.get('rx') or 0)
    if w <= 0 or h <= 0: return None
    rx = min(rx, w/2); ry = min(ry, h/2)
    if rx == 0 and ry == 0:
        return f'M {_f(x)},{_f(y)} L {_f(x+w)},{_f(y)} L {_f(x+w)},{_f(y+h)} L {_f(x)},{_f(y+h)} Z'
    kx, ky = KAPPA*rx, KAPPA*ry
    return (f'M {_f(x+rx)},{_f(y)} L {_f(x+w-rx)},{_f(y)} '
            f'C {_f(x+w-rx+kx)},{_f(y)} {_f(x+w)},{_f(y+ry-ky)} {_f(x+w)},{_f(y+ry)} '
            f'L {_f(x+w)},{_f(y+h-ry)} '
            f'C {_f(x+w)},{_f(y+h-ry+ky)} {_f(x+w-rx+kx)},{_f(y+h)} {_f(x+w-rx)},{_f(y+h)} '
            f'L {_f(x+rx)},{_f(y+h)} '
            f'C {_f(x+rx-kx)},{_f(y+h)} {_f(x)},{_f(y+h-ry+ky)} {_f(x)},{_f(y+h-ry)} '
            f'L {_f(x)},{_f(y+ry)} '
            f'C {_f(x)},{_f(y+ry-ky)} {_f(x+rx-kx)},{_f(y)} {_f(x+rx)},{_f(y)} Z')

def _ellipse_d(cx, cy, rx, ry):
    kx, ky = KAPPA*rx, KAPPA*ry
    return (f'M {_f(cx)},{_f(cy-ry)} '
            f'C {_f(cx+kx)},{_f(cy-ry)} {_f(cx+rx)},{_f(cy-ky)} {_f(cx+rx)},{_f(cy)} '
            f'C {_f(cx+rx)},{_f(cy+ky)} {_f(cx+kx)},{_f(cy+ry)} {_f(cx)},{_f(cy+ry)} '
            f'C {_f(cx-kx)},{_f(cy+ry)} {_f(cx-rx)},{_f(cy+ky)} {_f(cx-rx)},{_f(cy)} '
            f'C {_f(cx-rx)},{_f(cy-ky)} {_f(cx-kx)},{_f(cy-ry)} {_f(cx)},{_f(cy-ry)} Z')

def circle_to_d(el):
    cx=float(el.get('cx',0)); cy=float(el.get('cy',0)); r=float(el.get('r',0))
    return _ellipse_d(cx,cy,r,r) if r>0 else None

def ellipse_el_to_d(el):
    cx=float(el.get('cx',0)); cy=float(el.get('cy',0))
    rx=float(el.get('rx',0)); ry=float(el.get('ry',0))
    return _ellipse_d(cx,cy,rx,ry) if rx>0 and ry>0 else None

def points_to_d(el, close=False):
    pts = re.findall(r'[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?', el.get('points',''))
    if len(pts) < 4: return None
    coords = [float(v) for v in pts]
    pairs = [(coords[i], coords[i+1]) for i in range(0, len(coords)-1, 2)]
    d = f'M {_f(pairs[0][0])},{_f(pairs[0][1])}'
    for x,y in pairs[1:]: d += f' L {_f(x)},{_f(y)}'
    return d + ' Z' if close else d

def line_to_d(el):
    return (f'M {_f(float(el.get("x1",0)))},{_f(float(el.get("y1",0)))} '
            f'L {_f(float(el.get("x2",0)))},{_f(float(el.get("y2",0)))}')

def element_to_d(el):
    tag = _tag(el)
    if tag == 'rect':     return rect_to_d(el)
    if tag == 'circle':   return circle_to_d(el)
    if tag == 'ellipse':  return ellipse_el_to_d(el)
    if tag == 'polygon':  return points_to_d(el, close=True)
    if tag == 'polyline': return points_to_d(el, close=False)
    if tag == 'line':     return line_to_d(el)
    if tag == 'path':     return el.get('d','') or None
    return None

# ── Arc → cubic bezier conversion ───────────────────────────────────────────

def _arc_to_cubics(x1, y1, rx, ry, phi_deg, fa, fs, x2, y2):
    """
    Convert one SVG arc segment to a list of cubic bezier tuples.
    Each tuple is (cp1x, cp1y, cp2x, cp2y, ex, ey).
    Algorithm: endpoint → center parameterization, split into ≤90° segments,
    approximate each with the standard α formula.
    """
    if rx == 0 or ry == 0:
        return [(x2, y2, x2, y2, x2, y2)]

    phi = math.radians(phi_deg)
    cp, sp = math.cos(phi), math.sin(phi)

    # Midpoint transform
    dx, dy = (x1 - x2) / 2, (y1 - y2) / 2
    x1p =  cp*dx + sp*dy
    y1p = -sp*dx + cp*dy

    rx, ry = abs(rx), abs(ry)
    # Scale radii if too small
    lam = (x1p/rx)**2 + (y1p/ry)**2
    if lam > 1:
        s = math.sqrt(lam)
        rx *= s; ry *= s

    rx2, ry2 = rx*rx, ry*ry
    x1p2, y1p2 = x1p*x1p, y1p*y1p

    # Center in rotated frame
    num = max(0, rx2*ry2 - rx2*y1p2 - ry2*x1p2)
    den = rx2*y1p2 + ry2*x1p2
    sq = math.sqrt(num/den) if den else 0
    if fa == fs:
        sq = -sq
    cxp =  sq * rx * y1p / ry
    cyp = -sq * ry * x1p / rx

    # Center in original frame
    cx = cp*cxp - sp*cyp + (x1+x2)/2
    cy = sp*cxp + cp*cyp + (y1+y2)/2

    def _angle(ux, uy, vx, vy):
        n = math.hypot(ux, uy) * math.hypot(vx, vy)
        if n == 0: return 0
        a = math.acos(max(-1, min(1, (ux*vx + uy*vy) / n)))
        return -a if ux*vy - uy*vx < 0 else a

    theta1 = _angle(1, 0, (x1p-cxp)/rx, (y1p-cyp)/ry)
    dtheta  = _angle((x1p-cxp)/rx, (y1p-cyp)/ry,
                     (-x1p-cxp)/rx, (-y1p-cyp)/ry)

    if not fs and dtheta > 0: dtheta -= 2*math.pi
    elif fs and dtheta < 0:   dtheta += 2*math.pi

    n_segs = max(1, math.ceil(abs(dtheta) / (math.pi/2)))
    d_seg  = dtheta / n_segs

    cubics = []
    t = theta1
    for _ in range(n_segs):
        # α coefficient for this segment
        tan_half = math.tan(d_seg/2)
        alpha = math.sin(d_seg) * (math.sqrt(4 + 3*tan_half*tan_half) - 1) / 3

        cos_t,  sin_t  = math.cos(t),        math.sin(t)
        cos_t2, sin_t2 = math.cos(t+d_seg),  math.sin(t+d_seg)

        # Derivative vectors (in original frame)
        def _deriv(ct, st):
            ddx = -rx*st; ddy = ry*ct
            return cp*ddx - sp*ddy, sp*ddx + cp*ddy

        # Endpoint coordinates
        def _pt_on_arc(ct, st):
            return cx + cp*rx*ct - sp*ry*st, cy + sp*rx*ct + cp*ry*st

        px1, py1 = _pt_on_arc(cos_t,  sin_t)
        px2, py2 = _pt_on_arc(cos_t2, sin_t2)
        d1x, d1y = _deriv(cos_t,  sin_t)
        d2x, d2y = _deriv(cos_t2, sin_t2)

        cubics.append((
            px1 + alpha*d1x, py1 + alpha*d1y,   # cp1
            px2 - alpha*d2x, py2 - alpha*d2y,   # cp2
            px2, py2                              # endpoint
        ))
        t += d_seg

    return cubics


def normalize_arcs(d):
    """Replace all A/a arc commands in a path d string with cubic beziers."""
    if 'A' not in d and 'a' not in d:
        return d

    segs = expand_path(d)
    parts = []
    cx = cy = mx = my = 0.0

    for cmd, args in segs:
        if cmd == 'M':
            cx, cy = mx, my = args[0], args[1]
            parts.append(f'M {_f(cx)},{_f(cy)}')
        elif cmd == 'm':
            cx, cy = cx+args[0], cy+args[1]; mx, my = cx, cy
            parts.append(f'M {_f(cx)},{_f(cy)}')
        elif cmd == 'L':
            cx, cy = args[0], args[1]
            parts.append(f'L {_f(cx)},{_f(cy)}')
        elif cmd == 'l':
            cx += args[0]; cy += args[1]
            parts.append(f'L {_f(cx)},{_f(cy)}')
        elif cmd == 'H':
            cx = args[0]; parts.append(f'L {_f(cx)},{_f(cy)}')
        elif cmd == 'h':
            cx += args[0]; parts.append(f'L {_f(cx)},{_f(cy)}')
        elif cmd == 'V':
            cy = args[0]; parts.append(f'L {_f(cx)},{_f(cy)}')
        elif cmd == 'v':
            cy += args[0]; parts.append(f'L {_f(cx)},{_f(cy)}')
        elif cmd == 'C':
            x1,y1,x2,y2,x,y = args
            parts.append(f'C {_f(x1)},{_f(y1)} {_f(x2)},{_f(y2)} {_f(x)},{_f(y)}')
            cx, cy = x, y
        elif cmd == 'c':
            x1,y1,x2,y2,dx,dy = args
            ax1,ay1=cx+x1,cy+y1; ax2,ay2=cx+x2,cy+y2; ax,ay=cx+dx,cy+dy
            parts.append(f'C {_f(ax1)},{_f(ay1)} {_f(ax2)},{_f(ay2)} {_f(ax)},{_f(ay)}')
            cx, cy = ax, ay
        elif cmd == 'S':
            x2,y2,x,y = args
            parts.append(f'S {_f(x2)},{_f(y2)} {_f(x)},{_f(y)}'); cx,cy=x,y
        elif cmd == 's':
            dx2,dy2,dx,dy = args
            parts.append(f's {_f(dx2)},{_f(dy2)} {_f(dx)},{_f(dy)}'); cx+=dx; cy+=dy
        elif cmd == 'Q':
            qx,qy,x,y = args
            parts.append(f'Q {_f(qx)},{_f(qy)} {_f(x)},{_f(y)}'); cx,cy=x,y
        elif cmd == 'q':
            dqx,dqy,dx,dy = args
            parts.append(f'q {_f(dqx)},{_f(dqy)} {_f(dx)},{_f(dy)}'); cx+=dx; cy+=dy
        elif cmd == 'T':
            x,y = args; parts.append(f'T {_f(x)},{_f(y)}'); cx,cy=x,y
        elif cmd == 't':
            dx,dy = args; parts.append(f't {_f(dx)},{_f(dy)}'); cx+=dx; cy+=dy
        elif cmd == 'A':
            rx,ry,phi,fa,fs,x,y = args
            for c in _arc_to_cubics(cx,cy,rx,ry,phi,int(fa),int(fs),x,y):
                parts.append(f'C {_f(c[0])},{_f(c[1])} {_f(c[2])},{_f(c[3])} {_f(c[4])},{_f(c[5])}')
            cx, cy = x, y
        elif cmd == 'a':
            rx,ry,phi,fa,fs,dx,dy = args
            x, y = cx+dx, cy+dy
            for c in _arc_to_cubics(cx,cy,rx,ry,phi,int(fa),int(fs),x,y):
                parts.append(f'C {_f(c[0])},{_f(c[1])} {_f(c[2])},{_f(c[3])} {_f(c[4])},{_f(c[5])}')
            cx, cy = x, y
        elif cmd in ('Z','z'):
            parts.append('Z'); cx, cy = mx, my

    return ' '.join(parts)

# ── Transform parsing ────────────────────────────────────────────────────────

def parse_transform(s):
    fns = []
    for m in re.finditer(r'(\w+)\(([^)]*)\)', s or ''):
        name = m.group(1)
        vals = [float(v) for v in re.split(r'[\s,]+', m.group(2).strip()) if v]
        if name == 'translate':
            tx=vals[0]; ty=vals[1] if len(vals)>1 else 0.0
            fns.append(lambda x,y,tx=tx,ty=ty:(x+tx,y+ty))
        elif name == 'scale':
            sx=vals[0]; sy=vals[1] if len(vals)>1 else sx
            fns.append(lambda x,y,sx=sx,sy=sy:(x*sx,y*sy))
        elif name == 'rotate':
            a=math.radians(vals[0])
            cx=vals[1] if len(vals)>1 else 0.0; cy=vals[2] if len(vals)>2 else 0.0
            ca,sa=math.cos(a),math.sin(a)
            fns.append(lambda x,y,cx=cx,cy=cy,ca=ca,sa=sa:
                ((x-cx)*ca-(y-cy)*sa+cx,(x-cx)*sa+(y-cy)*ca+cy))
        elif name == 'matrix':
            a,b,c,d,e,f=vals[:6]
            fns.append(lambda x,y,a=a,b=b,c=c,d=d,e=e,f=f:(a*x+c*y+e,b*x+d*y+f))
    if not fns: return None
    def composed(x,y):
        for fn in reversed(fns): x,y = fn(x,y)
        return x,y
    return composed

def chain(outer, inner):
    if outer is None: return inner
    if inner is None: return outer
    return lambda x,y: outer(*inner(x,y))

# ── SVG path tokenizer ───────────────────────────────────────────────────────

_CMD_ARGC = dict(M=2,m=2,L=2,l=2,H=1,h=1,V=1,v=1,
                 C=6,c=6,S=4,s=4,Q=4,q=4,T=2,t=2,A=7,a=7,Z=0,z=0)

def expand_path(d):
    toks = re.findall(
        r'[MLHVCSQTAZmlhvcsqtaz]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?', d)
    segs, cmd, buf = [], None, []
    for t in toks:
        if t.isalpha():
            if cmd is not None: segs.append((cmd, buf))
            cmd, buf = t, []
        else:
            buf.append(float(t))
    if cmd is not None: segs.append((cmd, buf))
    out = []
    for cmd, args in segs:
        n = _CMD_ARGC.get(cmd.upper(), 0)
        if n == 0:
            out.append((cmd, [])); continue
        for i in range(0, max(len(args), n), n):
            chunk = args[i:i+n]
            if len(chunk) < n: break
            c = ('L' if cmd=='M' else 'l' if cmd=='m' else cmd) if i>0 else cmd
            out.append((c, chunk))
    return out

# ── Stroke expansion ─────────────────────────────────────────────────────────

def _norm(dx, dy):
    L = math.hypot(dx, dy)
    return (dx/L, dy/L) if L > 1e-10 else (0.0, 0.0)

def _left_normal(dx, dy):
    nx, ny = _norm(dx, dy)
    return -ny, nx

def _miter_offset(n1, n2, hw, limit=4.0):
    """Offset vector at a vertex join for half-width hw."""
    sx, sy = n1[0]+n2[0], n1[1]+n2[1]
    denom = n1[0]*sx + n1[1]*sy
    if abs(denom) < 1e-8:
        return n1[0]*hw, n1[1]*hw
    scale = hw / denom
    if abs(scale) > limit:
        # Bevel fallback
        return n1[0]*hw, n1[1]*hw
    return sx*scale, sy*scale

def _bezier3(p0, p1, p2, p3, t):
    mt = 1-t
    return (mt**3*p0[0]+3*mt**2*t*p1[0]+3*mt*t**2*p2[0]+t**3*p3[0],
            mt**3*p0[1]+3*mt**2*t*p1[1]+3*mt*t**2*p2[1]+t**3*p3[1])

def _sample_segments(d, n=20):
    """Yield (points_list, is_closed) polylines by sampling the path."""
    segs = expand_path(d)
    current = []
    mx = my = cx = cy = 0.0
    last_ctrl = None
    polylines = []

    def ap(x, y):
        if not current or current[-1] != (x, y):
            current.append((x, y))

    def bez(p0, p1, p2, p3):
        for i in range(1, n+1):
            ap(*_bezier3(p0, p1, p2, p3, i/n))

    for cmd, args in segs:
        if cmd not in ('C','c','S','s','Q','q','T','t'):
            last_ctrl = None
        if cmd == 'M':
            if len(current) >= 2: polylines.append((list(current), False))
            current.clear(); cx,cy = mx,my = args[0],args[1]; ap(cx,cy)
        elif cmd == 'm':
            if len(current) >= 2: polylines.append((list(current), False))
            current.clear(); cx,cy = cx+args[0],cy+args[1]; mx,my = cx,cy; ap(cx,cy)
        elif cmd == 'L':  cx,cy = args; ap(cx,cy)
        elif cmd == 'l':  cx+=args[0]; cy+=args[1]; ap(cx,cy)
        elif cmd == 'H':  cx=args[0]; ap(cx,cy)
        elif cmd == 'h':  cx+=args[0]; ap(cx,cy)
        elif cmd == 'V':  cy=args[0]; ap(cx,cy)
        elif cmd == 'v':  cy+=args[0]; ap(cx,cy)
        elif cmd == 'C':
            x1,y1,x2,y2,x,y=args; bez((cx,cy),(x1,y1),(x2,y2),(x,y))
            last_ctrl=(x2,y2); cx,cy=x,y
        elif cmd == 'c':
            x1,y1,x2,y2,dx,dy=args
            p1=(cx+x1,cy+y1); p2=(cx+x2,cy+y2); p3=(cx+dx,cy+dy)
            bez((cx,cy),p1,p2,p3); last_ctrl=p2; cx,cy=p3
        elif cmd == 'S':
            x2,y2,x,y=args; lx,ly=last_ctrl if last_ctrl else (cx,cy)
            bez((cx,cy),(2*cx-lx,2*cy-ly),(x2,y2),(x,y)); last_ctrl=(x2,y2); cx,cy=x,y
        elif cmd == 's':
            dx2,dy2,dx,dy=args; lx,ly=last_ctrl if last_ctrl else (cx,cy)
            p2=(cx+dx2,cy+dy2); p3=(cx+dx,cy+dy)
            bez((cx,cy),(2*cx-lx,2*cy-ly),p2,p3); last_ctrl=p2; cx,cy=p3
        elif cmd == 'Q':
            qx,qy,x,y=args
            p1=(cx+2/3*(qx-cx),cy+2/3*(qy-cy)); p2=(x+2/3*(qx-x),y+2/3*(qy-y))
            bez((cx,cy),p1,p2,(x,y)); last_ctrl=(qx,qy); cx,cy=x,y
        elif cmd == 'q':
            dqx,dqy,dx,dy=args; qx,qy=cx+dqx,cy+dqy; x,y=cx+dx,cy+dy
            p1=(cx+2/3*(qx-cx),cy+2/3*(qy-cy)); p2=(x+2/3*(qx-x),y+2/3*(qy-y))
            bez((cx,cy),p1,p2,(x,y)); last_ctrl=(qx,qy); cx,cy=x,y
        elif cmd in ('Z','z'):
            if len(current) >= 2: polylines.append((list(current), True))
            current.clear(); current.append((mx,my)); cx,cy=mx,my
        elif cmd == 'A':  cx,cy=args[5],args[6]; ap(cx,cy)
        elif cmd == 'a':  cx+=args[5]; cy+=args[6]; ap(cx,cy)

    if len(current) >= 2:
        polylines.append((list(current), False))
    return polylines

def _offset_points(points, hw, closed):
    """Compute left (+hw) and right (-hw) offset point lists."""
    n = len(points)
    if n < 2: return [], []

    seg_normals = []
    for i in range(n-1):
        seg_normals.append(_left_normal(points[i+1][0]-points[i][0],
                                        points[i+1][1]-points[i][1]))
    if closed:
        seg_normals.append(_left_normal(points[0][0]-points[-1][0],
                                        points[0][1]-points[-1][1]))

    left_pts, right_pts = [], []
    for i in range(n):
        if closed:
            n1 = seg_normals[(i-1) % n]
            n2 = seg_normals[i % n]
        elif i == 0:   n1 = n2 = seg_normals[0]
        elif i == n-1: n1 = n2 = seg_normals[-1]
        else:          n1, n2 = seg_normals[i-1], seg_normals[i]
        ox, oy = _miter_offset(n1, n2, hw)
        left_pts.append((points[i][0]+ox, points[i][1]+oy))
        right_pts.append((points[i][0]-ox, points[i][1]-oy))
    return left_pts, right_pts

def stroke_to_outline(d, stroke_width):
    """
    Expand a stroked path to a filled path representing the stroke area.

    Open paths:   left side forward + right side backward + butt caps → one closed ring
    Closed paths: outer ring (CW) + inner ring (CCW) as two subpaths;
                  with nonzero fill rule this renders as a donut.
    """
    hw = stroke_width / 2.0
    parts = []

    for points, closed in _sample_segments(d):
        if len(points) < 2: continue
        lp, rp = _offset_points(points, hw, closed)
        if not lp: continue

        if closed:
            # Outer ring forward (same winding as original = CW for typical SVG shapes)
            out  = f'M {_f(*lp[0])}'
            for p in lp[1:]: out += f' L {_f(*p)}'
            out += ' Z'
            # Inner ring reversed → opposite winding (CCW) → creates donut with nonzero rule
            out += f' M {_f(*rp[-1])}'
            for p in reversed(rp[:-1]): out += f' L {_f(*p)}'
            out += ' Z'
        else:
            # Open path: left forward, butt cap at end, right backward, butt cap at start
            out  = f'M {_f(*lp[0])}'
            for p in lp[1:]: out += f' L {_f(*p)}'
            out += f' L {_f(*rp[-1])}'
            for p in reversed(rp[:-1]): out += f' L {_f(*p)}'
            out += ' Z'
        parts.append(out)

    return ' '.join(parts) if parts else d

# ── SVG tree traversal ───────────────────────────────────────────────────────

SHAPE_TAGS = {'path','rect','circle','ellipse','polygon','polyline','line'}

def _is_background_rect(el, vb_x, vb_y, vb_w, vb_h):
    """True if this rect covers the full viewBox — a canvas background, not content."""
    if _tag(el) != 'rect':
        return False
    try:
        x = float(el.get('x', 0)); y = float(el.get('y', 0))
        w = float(el.get('width', 0)); h = float(el.get('height', 0))
    except ValueError:
        return False
    tol = 0.02  # 2% tolerance
    return (abs(x - vb_x) <= vb_w * tol and
            abs(y - vb_y) <= vb_h * tol and
            abs(w - vb_w) <= vb_w * tol and
            abs(h - vb_h) <= vb_h * tol)

def collect(el, acc_xfm=None, inh_fill='000000', gradients=None, vb=None):
    """Yield (d, transform_fn, fill_hex) for every shape in the SVG tree."""
    if gradients is None: gradients = {}

    local_xfm = parse_transform(el.get('transform',''))
    xfm = chain(acc_xfm, local_xfm)
    ident = lambda x,y: (x,y)

    fill_hex, stroke_hex, stroke_width = resolve_fill_stroke(el, gradients, inh_fill)
    next_inh = fill_hex if fill_hex is not None else inh_fill

    tag = _tag(el)
    if tag in SHAPE_TAGS:
        # Skip canvas background rectangles
        if vb and _is_background_rect(el, *vb):
            pass
        else:
            d = element_to_d(el)
            if d: d = normalize_arcs(d)
            if d:
                if fill_hex is not None:
                    yield (d, xfm or ident, fill_hex)
                if stroke_hex and stroke_width > 0:
                    outlined = stroke_to_outline(d, stroke_width)
                    yield (outlined, xfm or ident, stroke_hex)
                elif stroke_hex and fill_hex is None:
                    yield (d, xfm or ident, stroke_hex)

    for child in el:
        yield from collect(child, xfm, next_inh, gradients, vb)

# ── DrawingML path builder ───────────────────────────────────────────────────

def _el(name):     return etree.Element(f'{{{_A}}}{name}')
def _sub(p, name): return etree.SubElement(p, f'{{{_A}}}{name}')
def _pt(x, y):
    e = _el('pt'); e.set('x', str(x)); e.set('y', str(y)); return e

def make_a_path(d, xfm, vb_x, vb_y, vb_w, vb_h):
    cw = COORD
    ch = round(COORD * vb_h / vb_w) if vb_w else COORD

    def sc(x, y):
        tx, ty = xfm(x, y)
        return round((tx-vb_x)/vb_w*cw), round((ty-vb_y)/vb_h*ch)

    a_path = _el('path')
    a_path.set('w', str(cw)); a_path.set('h', str(ch))

    cx = cy = mx = my = 0.0
    last_ctrl = None

    for cmd, args in expand_path(d):
        if cmd not in ('C','c','S','s','Q','q','T','t'): last_ctrl = None

        if   cmd=='M': cx,cy=mx,my=args;  _sub(a_path,'moveTo').append(_pt(*sc(cx,cy)))
        elif cmd=='m': cx,cy=mx,my=cx+args[0],cy+args[1]; _sub(a_path,'moveTo').append(_pt(*sc(cx,cy)))
        elif cmd=='L': cx,cy=args;         _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='l': cx+=args[0];cy+=args[1]; _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='H': cx=args[0];         _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='h': cx+=args[0];        _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='V': cy=args[0];         _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='v': cy+=args[0];        _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='C':
            x1,y1,x2,y2,x,y=args; e=_sub(a_path,'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl=(x2,y2); cx,cy=x,y
        elif cmd=='c':
            x1,y1,x2,y2,dx,dy=args
            ax1,ay1=cx+x1,cy+y1; ax2,ay2=cx+x2,cy+y2; ax,ay=cx+dx,cy+dy
            e=_sub(a_path,'cubicBezTo')
            e.append(_pt(*sc(ax1,ay1))); e.append(_pt(*sc(ax2,ay2))); e.append(_pt(*sc(ax,ay)))
            last_ctrl=(ax2,ay2); cx,cy=ax,ay
        elif cmd=='S':
            x2,y2,x,y=args; lx,ly=last_ctrl if last_ctrl else (cx,cy)
            x1,y1=2*cx-lx,2*cy-ly; e=_sub(a_path,'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl=(x2,y2); cx,cy=x,y
        elif cmd=='s':
            dx2,dy2,dx,dy=args; x2,y2=cx+dx2,cy+dy2; x,y=cx+dx,cy+dy
            lx,ly=last_ctrl if last_ctrl else (cx,cy); x1,y1=2*cx-lx,2*cy-ly
            e=_sub(a_path,'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl=(x2,y2); cx,cy=x,y
        elif cmd=='Q':
            qx,qy,x,y=args
            x1,y1=cx+2/3*(qx-cx),cy+2/3*(qy-cy); x2,y2=x+2/3*(qx-x),y+2/3*(qy-y)
            e=_sub(a_path,'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl=(qx,qy); cx,cy=x,y
        elif cmd=='q':
            dqx,dqy,dx,dy=args; qx,qy=cx+dqx,cy+dqy; x,y=cx+dx,cy+dy
            x1,y1=cx+2/3*(qx-cx),cy+2/3*(qy-cy); x2,y2=x+2/3*(qx-x),y+2/3*(qy-y)
            e=_sub(a_path,'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl=(qx,qy); cx,cy=x,y
        elif cmd in ('Z','z'): _sub(a_path,'close'); cx,cy=mx,my
        elif cmd=='A':  cx,cy=args[5],args[6]; _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))
        elif cmd=='a':  cx+=args[5];cy+=args[6]; _sub(a_path,'lnTo').append(_pt(*sc(cx,cy)))

    return a_path

# ── PPTX assembly ────────────────────────────────────────────────────────────

def add_shapes(slide, paths, vb_x, vb_y, vb_w, vb_h):
    aspect = vb_w/vb_h if vb_h else 1
    if aspect >= SLIDE_W/SLIDE_H:
        sw = round(SLIDE_W*0.85); sh = round(sw/aspect)
    else:
        sh = round(SLIDE_H*0.85); sw = round(sh*aspect)
    ox = (SLIDE_W-sw)//2; oy = (SLIDE_H-sh)//2

    cw = COORD
    ch = round(COORD*vb_h/vb_w) if vb_w else COORD

    for sp_id, (d, xfm, fill_hex) in enumerate(paths, start=2):
        a_path = make_a_path(d, xfm, vb_x, vb_y, vb_w, vb_h)
        pts = list(a_path.iter(f'{{{_A}}}pt'))
        if pts:
            xs=[int(pt.get('x')) for pt in pts]; ys=[int(pt.get('y')) for pt in pts]
            px_min,py_min=min(xs),min(ys); px_max,py_max=max(xs),max(ys)
            pw=max(px_max-px_min,1); ph=max(py_max-py_min,1)
            for pt in pts:
                pt.set('x',str(int(pt.get('x'))-px_min))
                pt.set('y',str(int(pt.get('y'))-py_min))
            a_path.set('w',str(pw)); a_path.set('h',str(ph))
            shape_ox=ox+round(px_min/cw*sw); shape_oy=oy+round(py_min/ch*sh)
            shape_sw=max(round(pw/cw*sw),1); shape_sh=max(round(ph/ch*sh),1)
        else:
            shape_ox,shape_oy,shape_sw,shape_sh = ox,oy,sw,sh

        hx = fill_hex.upper().lstrip('#')
        sp = etree.fromstring(
            f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{sp_id}" name="Shape {sp_id-1}"/>'
            f'<p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr>'
            f'<p:nvPr/></p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{shape_ox}" y="{shape_oy}"/>'
            f'<a:ext cx="{shape_sw}" cy="{shape_sh}"/></a:xfrm>'
            f'<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
            f'<a:rect l="0" t="0" r="r" b="b"/><a:pathLst/></a:custGeom>'
            f'<a:solidFill><a:srgbClr val="{hx}"/></a:solidFill>'
            f'<a:ln><a:noFill/></a:ln>'
            f'</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
            f'</p:sp>'
        )
        sp.find(f'.//{{{_A}}}pathLst').append(a_path)
        slide.shapes._spTree.append(sp)

# ── Entry point ──────────────────────────────────────────────────────────────

def convert(svg_files, out_dir=None):
    """Convert one or more SVG files into a single PPTX — one slide per SVG.

    Returns the Path of the written PPTX, or None on failure.
    out_dir: directory for output file; defaults to alongside the first input.
    """
    svg_files = [f for f in svg_files if Path(f).suffix.lower() == '.svg']
    if not svg_files:
        print('No .svg files provided.', file=sys.stderr)
        return None

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for svg_file in svg_files:
        p = Path(svg_file)
        root = ET.parse(p).getroot()
        vb = root.get('viewBox','0 0 100 100')
        vb_x,vb_y,vb_w,vb_h = [float(v) for v in re.split(r'[\s,]+', vb.strip())]
        gradients = extract_gradients(root)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = list(collect(root, gradients=gradients, vb=(vb_x, vb_y, vb_w, vb_h)))
        add_shapes(slide, shapes, vb_x, vb_y, vb_w, vb_h)
        print(f'  + {p.name} ({len(shapes)} shapes)')

    first = Path(svg_files[0])
    stem = first.stem if len(svg_files)==1 else 'slides'
    dest_dir = Path(out_dir) if out_dir else first.parent
    base = dest_dir / f'{stem}.pptx'
    out = base; n = 2
    while out.exists():
        out = dest_dir / f'{stem} {n}.pptx'; n += 1
    prs.save(str(out))
    print(f'✓ {out}')
    return out

if __name__ == '__main__':
    if len(sys.argv) < 2:
        sys.exit(f'Usage: {Path(sys.argv[0]).name} file.svg [...]')
    try:
        convert(sys.argv[1:])
    except Exception as ex:
        print(f'✗ {ex}', file=sys.stderr)
        sys.exit(1)
